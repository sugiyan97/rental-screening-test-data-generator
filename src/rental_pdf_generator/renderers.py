"""HTMLテンプレートを各出力形式に変換するレンダラー群。

Playwright の Page に描画済みの HTML から、
- pdf / png / jpg: Playwright の機能でそのまま出力
- xlsx / docx / csv / pptx: DOM から見出し・段落・表を抽出して再構成
という2系統で出力する。
"""

import csv
from collections.abc import Callable
from pathlib import Path
from typing import Any

OUTPUT_FORMATS: tuple[str, ...] = ("pdf", "png", "jpg", "xlsx", "docx", "csv", "pptx")

# 出力形式ごとの出力先サブディレクトリ（ケースディレクトリ直下）
SUBDIR_BY_FORMAT: dict[str, str] = {fmt: fmt for fmt in OUTPUT_FORMATS}

# 出力形式ごとの拡張子
EXTENSION_BY_FORMAT: dict[str, str] = {fmt: fmt for fmt in OUTPUT_FORMATS}


class UnsupportedOutputFormatError(Exception):
    pass


class PdfPasswordNotSupportedError(Exception):
    """pdf 以外の出力形式に pdf_password が指定された場合のエラー。"""


# DOM から「見出し / 段落 / 表」のブロック列を抽出する。
# ブロック要素を含まない最小単位のボックスを1段落として扱うため、
# ラベルと値を横並びにしたレイアウトでも1行にまとまる。
_EXTRACT_BLOCKS_JS = r"""
() => {
  const BOXY = new Set([
    'block', 'flex', 'grid', 'list-item', 'flow-root',
    'table', 'table-row', 'table-cell', 'table-row-group',
  ]);
  const hidden = (el) => {
    const s = getComputedStyle(el);
    return s.display === 'none' || s.visibility === 'hidden';
  };
  const isBoxy = (el) => BOXY.has(getComputedStyle(el).display);
  const hasBoxyDescendant = (el) => {
    for (const child of el.children) {
      if (hidden(child)) continue;
      if (isBoxy(child) || hasBoxyDescendant(child)) return true;
    }
    return false;
  };
  const clean = (value) => (value || '').replace(/\s+/g, ' ').trim();
  const textOf = (el) => clean(el.innerText || el.textContent);
  const tableRows = (table) => {
    const rows = [];
    for (const tr of table.querySelectorAll('tr')) {
      if (hidden(tr)) continue;
      const cells = [];
      for (const cell of tr.children) {
        if (hidden(cell)) continue;
        cells.push(textOf(cell));
      }
      if (cells.some((c) => c !== '')) rows.push(cells);
    }
    return rows;
  };
  const blocks = [];
  const walk = (el) => {
    for (const node of el.childNodes) {
      if (node.nodeType === Node.TEXT_NODE) {
        const text = clean(node.textContent);
        if (text) blocks.push({ type: 'text', text });
        continue;
      }
      if (node.nodeType !== Node.ELEMENT_NODE) continue;
      const tag = node.tagName.toLowerCase();
      if (tag === 'script' || tag === 'style' || hidden(node)) continue;
      if (tag === 'table') {
        const rows = tableRows(node);
        if (rows.length) blocks.push({ type: 'table', rows });
        continue;
      }
      if (/^h[1-6]$/.test(tag)) {
        const text = textOf(node);
        if (text) blocks.push({ type: 'heading', level: Number(tag[1]), text });
        continue;
      }
      if (!hasBoxyDescendant(node)) {
        const text = textOf(node);
        if (text) blocks.push({ type: 'text', text });
        continue;
      }
      walk(node);
    }
  };
  walk(document.body);
  return blocks;
}
"""

_JP_FONT = "Yu Gothic"


def extract_blocks(page) -> list[dict[str, Any]]:
    """描画済みページから見出し・段落・表のブロック列を取り出す。"""
    return page.evaluate(_EXTRACT_BLOCKS_JS)


def _render_pdf(page, path: Path, title: str) -> None:
    page.pdf(path=str(path), format="A4", print_background=True)


def _render_png(page, path: Path, title: str) -> None:
    page.screenshot(path=str(path), type="png", full_page=True)


def _render_jpg(page, path: Path, title: str) -> None:
    page.screenshot(path=str(path), type="jpeg", quality=90, full_page=True)


def _render_docx(page, path: Path, title: str) -> None:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.shared import Pt

    document = Document()
    style = document.styles["Normal"]
    style.font.name = _JP_FONT
    style.font.size = Pt(10.5)
    style.element.get_or_add_rPr().get_or_add_rFonts().set(qn("w:eastAsia"), _JP_FONT)

    for block in extract_blocks(page):
        if block["type"] == "heading":
            document.add_heading(block["text"], level=min(block.get("level", 1), 4))
        elif block["type"] == "text":
            document.add_paragraph(block["text"])
        else:
            rows = block["rows"]
            table = document.add_table(rows=len(rows), cols=max(len(r) for r in rows))
            table.style = "Table Grid"
            for row_idx, row in enumerate(rows):
                for col_idx, value in enumerate(row):
                    table.cell(row_idx, col_idx).text = value
            document.add_paragraph("")

    document.save(str(path))


def _render_xlsx(page, path: Path, title: str) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = title[:31] or "document"
    sheet.column_dimensions["A"].width = 40

    row_idx = 1
    for block in extract_blocks(page):
        if block["type"] in ("heading", "text"):
            cell = sheet.cell(row=row_idx, column=1, value=block["text"])
            if block["type"] == "heading":
                cell.font = Font(bold=True)
            row_idx += 1
        else:
            for row in block["rows"]:
                for col_idx, value in enumerate(row, start=1):
                    sheet.cell(row=row_idx, column=col_idx, value=value)
                row_idx += 1
            row_idx += 1

    workbook.save(str(path))


def _render_csv(page, path: Path, title: str) -> None:
    with path.open("w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        for block in extract_blocks(page):
            if block["type"] in ("heading", "text"):
                writer.writerow([block["text"]])
            else:
                writer.writerows(block["rows"])
                writer.writerow([])


_PPTX_LINES_PER_SLIDE = 12


def _render_pptx(page, path: Path, title: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]

    def add_slide(heading: str) -> Any:
        slide = presentation.slides.add_slide(blank_layout)
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = heading
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.name = _JP_FONT
        return slide

    def add_lines(slide, lines: list[str]) -> None:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.5))
        frame = box.text_frame
        frame.word_wrap = True
        for i, line in enumerate(lines):
            paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            run = paragraph.add_run()
            run.text = line
            run.font.size = Pt(12)
            run.font.name = _JP_FONT

    heading = title
    pending: list[str] = []
    slide_count = 0

    def flush() -> None:
        nonlocal pending, slide_count
        if not pending:
            return
        for start in range(0, len(pending), _PPTX_LINES_PER_SLIDE):
            add_lines(add_slide(heading), pending[start : start + _PPTX_LINES_PER_SLIDE])
            slide_count += 1
        pending = []

    for block in extract_blocks(page):
        if block["type"] == "heading":
            flush()
            heading = block["text"]
        elif block["type"] == "text":
            pending.append(block["text"])
        else:
            flush()
            rows = block["rows"]
            slide = add_slide(heading)
            shape = slide.shapes.add_table(
                len(rows),
                max(len(r) for r in rows),
                Inches(0.5),
                Inches(1.2),
                Inches(9.0),
                Inches(0.3 * len(rows)),
            )
            for row_idx, row in enumerate(rows):
                for col_idx, value in enumerate(row):
                    cell = shape.table.cell(row_idx, col_idx)
                    cell.text = value
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
                            run.font.name = _JP_FONT
            slide_count += 1

    flush()
    if slide_count == 0:
        add_slide(heading)

    presentation.save(str(path))


_RENDERERS: dict[str, Callable[[Any, Path, str], None]] = {
    "pdf": _render_pdf,
    "png": _render_png,
    "jpg": _render_jpg,
    "xlsx": _render_xlsx,
    "docx": _render_docx,
    "csv": _render_csv,
    "pptx": _render_pptx,
}


def encrypt_pdf(path: Path, password: str) -> None:
    """生成済みの PDF を user / owner 両方のパスワードで暗号化する（上書き保存）。

    「開けない PDF」をテストデータとして作るための処理。パスワードなしでは
    pikepdf / 一般的な PDF ビューアで開けなくなる。
    """
    import pikepdf

    with pikepdf.open(path, allow_overwriting_input=True) as pdf:
        pdf.save(
            path,
            encryption=pikepdf.Encryption(user=password, owner=password, R=6),
        )


def render_document(page, output_format: str, path: Path, title: str = "") -> None:
    """描画済みの page を output_format に応じて path へ書き出す。"""
    renderer = _RENDERERS.get(output_format)
    if renderer is None:
        raise UnsupportedOutputFormatError(
            f"未対応の出力形式です。\n"
            f"  output_format: {output_format}\n"
            f"  対応済み: {list(_RENDERERS.keys())}"
        )
    path.parent.mkdir(parents=True, exist_ok=True)
    renderer(page, path, title)
